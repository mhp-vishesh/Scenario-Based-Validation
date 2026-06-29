"""Build the 6-slide Scenario-Based Validation story on the real MHP template.

This is the consulting storyline placed on the genuine MHP Company Profile
master, so it inherits the actual MHP logo, the copyright footer, the page
numbers, and the Segoe UI house styling. Content is drawn into the master's
body area (title top band, body from y 1.6, footer untouched).

Output: Scenario-Based_Validation_MHP_Deck.pptx (workspace root)
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "MHP Company Profile EN.pptx"
SUT = ROOT / "POC" / "outputs" / "sut"
OUT = ROOT / "Scenario-Based_Validation_MHP_Deck.pptx"

YOLO_FOG = SUT / "20151221120048-D6-AGGRESSIVE-MOTORWAY_fog_dusk_pedestrian_occluded_emergence_000_yolo.jpg"
YOLO_RAIN = SUT / "20151221120048-D6-AGGRESSIVE-MOTORWAY_rain_night_stalled_vehicle_sudden_braking_001_yolo.jpg"
YOLO_SUN = SUT / "20151221120048-D6-AGGRESSIVE-MOTORWAY_clear_low_sun_glare_cyclist_jaywalk_cut_in_002_yolo.jpg"

# ------------------------------------------------------------------ palette (MHP theme)
BLUE = RGBColor(0x00, 0x00, 0x99)      # MHP primary
BLUE_T = RGBColor(0xEE, 0xF1, 0xFB)    # blue tint panel
INK = RGBColor(0x20, 0x24, 0x2B)       # primary text
SUB = RGBColor(0x5B, 0x64, 0x70)       # secondary text
HAIR = RGBColor(0xDC, 0xE1, 0xEA)      # hairline
GRID_OFF = RGBColor(0xE6, 0xE9, 0xEF)  # empty grid cell
LIME = RGBColor(0xDD, 0xEF, 0x03)      # MHP accent (sparing)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
R5 = RGBColor(0xC0, 0x39, 0x2B)
R4 = RGBColor(0xD9, 0x82, 0x2B)
R3 = RGBColor(0xB8, 0x86, 0x0B)

FONT = "Segoe UI"
FONT_L = "Segoe UI Light"
FONT_SB = "Segoe UI Semibold"

PW, PH = 13.333, 7.5
ML = 0.45              # master content left margin
CR = 12.88             # content right edge (clear of nothing; full bleed band)
TITLE_W = 10.99        # title width that clears the top-right logo
CW = CR - ML           # full content width
TOP = 1.5              # body content top (just under the accent rule)

prs = Presentation(str(TEMPLATE))
L = {x.name: x for x in prs.slide_layouts}


# ------------------------------------------------------------------ housekeeping
def delete_all_slides():
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        prs.part.drop_rel(sid.get(qn("r:id")))
        lst.remove(sid)


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(idx)


def remove_ph(slide, idx):
    ph = get_ph(slide, idx)
    geom = (Emu(ph.left).inches, Emu(ph.top).inches, Emu(ph.width).inches, Emu(ph.height).inches)
    ph._element.getparent().remove(ph._element)
    return geom


# ------------------------------------------------------------------ primitives
def _spc(run, pts):
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def rect(s, l, t, w, h, fill=None, line=None, lw=0.75):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def rrect(s, l, t, w, h, fill=None, line=None, lw=0.75, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def hline(s, x1, x2, y, color=HAIR):
    return rect(s, x1, y, x2 - x1, 0.012, fill=color)


def oval(s, l, t, w, h, fill=None, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, segs, size=13, color=INK, align=PP_ALIGN.LEFT, bold=False,
         font=FONT, before=0, after=4, line=1.06, spc=None, first=False):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = line
    if isinstance(segs, str):
        segs = [(segs, bold, color, font)]
    for seg in segs:
        text = seg[0]
        b = seg[1] if len(seg) > 1 else bold
        c = seg[2] if len(seg) > 2 else color
        f = seg[3] if len(seg) > 3 else font
        r = p.add_run()
        r.text = text
        r.font.name = f
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
        if spc is not None:
            _spc(r, spc)
    return p


def chevron(s, cx, cy):
    tf = tb(s, cx - 0.2, cy - 0.25, 0.4, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "\u203a", size=26, color=BLUE, bold=True, align=PP_ALIGN.CENTER, after=0, first=True)


def contain(path, bl, bt, bw, bh):
    iw, ih = Image.open(str(path)).size
    ar, bar = iw / ih, bw / bh
    if ar > bar:
        w, h = bw, bw / ar
    else:
        h, w = bh, bh * ar
    return bl + (bw - w) / 2, bt + (bh - h) / 2, w, h


def picture(s, path, bl, bt, bw, bh):
    l, t, w, h = contain(path, bl, bt, bw, bh)
    pic = s.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
    pic.line.color.rgb = HAIR
    pic.line.width = Pt(0.5)
    return l, t, w, h


def risk_pill(s, score, color, right_in, top_in):
    w, h = 0.92, 0.30
    shp = rrect(s, right_in - w - 0.08, top_in + 0.08, w, h, fill=color, radius=0.5)
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"Risk {score}/5"
    r.font.name = FONT_SB
    r.font.size = Pt(10.5)
    r.font.bold = True
    r.font.color.rgb = WHITE


# ------------------------------------------------------------------ header (MHP house style)
def header(s, eyebrow_text, title_text):
    """Section kicker top-left, message title in the master title band, accent rule."""
    # eyebrow above the title band, clear of the top-right logo
    tfe = tb(s, ML, 0.18, TITLE_W, 0.26)
    para(tfe, eyebrow_text.upper(), size=10.5, color=BLUE, bold=True, font=FONT_SB,
         spc=2.2, after=0, first=True)
    # reuse the master title placeholder so the title sits in the genuine band
    ph = get_ph(s, 0)
    ph.left, ph.top, ph.width, ph.height = Inches(ML), Inches(0.5), Inches(TITLE_W), Inches(0.66)
    tf = ph.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = title_text
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(23)
        r.font.bold = True
        r.font.color.rgb = INK
    # accent rule
    rect(s, ML, 1.24, 0.62, 0.05, fill=BLUE)
    rect(s, ML + 0.7, 1.24, 0.14, 0.05, fill=LIME)


# ================================================================== Slide 1: cover
def s_cover():
    s = prs.slides.add_slide(L["Title Slide"])
    # title (template cover; inherits MHP cover treatment + logo)
    ph = get_ph(s, 0)
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Scenario-Based Validation for ADAS"
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(38)
        r.font.bold = True
    # replace the narrow lime "Category" subtitle with a clean descriptive line
    geom = remove_ph(s, 1)
    sub = tb(s, geom[0], 4.18, 11.8, 1.5)
    para(sub, "Finding the failures a test fleet rarely sees, using synthetic scenarios "
              "and an automated safety judge.", size=18, color=WHITE, line=1.2, after=10, first=True)
    para(sub, "Proof of concept on NVIDIA Cosmos and AWS.", size=13.5, color=WHITE,
         bold=True, font=FONT_SB, after=0)


# ================================================================== Slide 2: executive summary
def s_exec():
    s = prs.slides.add_slide(L["Title only"])
    header(s, "Executive summary", "From a few clips to audited safety coverage")

    rows = [
        ("SITUATION", "OEMs must show ADAS behaves safely in rare conditions, and those "
                      "conditions are the hardest to capture on the road."),
        ("COMPLICATION", "Road data campaigns are slow and costly, yet they still miss most of "
                         "the hazardous edge cases that SOTIF asks you to address."),
        ("APPROACH", "We expand a small set of seed clips into many edge cases with NVIDIA "
                     "Cosmos, then score every clip with a judge that explains its reasoning."),
        ("RESULT", "The first run flagged a missed stalled vehicle at night and a late reaction "
                   "to an occluded pedestrian, each with a risk score and a written rationale."),
    ]
    x, w = ML, 7.05
    y, rh = TOP + 0.05, 1.2
    for label, body in rows:
        rect(s, x, y + 0.06, 0.045, 0.82, fill=BLUE)
        tf = tb(s, x + 0.22, y, w - 0.22, rh)
        para(tf, label, size=10.5, color=BLUE, bold=True, font=FONT_SB, spc=1.6, after=3, first=True)
        para(tf, body, size=12.5, color=INK, line=1.14, after=0)
        y += rh

    px = 8.05
    pw = CR - px
    rect(s, px, TOP + 0.05, pw, 4.85, fill=BLUE_T)
    rect(s, px, TOP + 0.05, pw, 0.06, fill=BLUE)
    tf = tb(s, px + 0.32, TOP + 0.42, pw - 0.64, 4.2)
    para(tf, "BOTTOM LINE", size=10.5, color=BLUE, bold=True, font=FONT_SB, spc=2, after=10, first=True)
    para(tf, "You build a defensible safety case for the conditions a fleet rarely sees, "
             "without launching a new data campaign.", size=16, color=INK, bold=True, line=1.2, after=14)
    para(tf, "Evidence is reproducible and standards-aligned for SOTIF (ISO 21448) and "
             "ISO 26262.", size=12, color=SUB, line=1.2, after=0)


# ================================================================== Slide 3: the challenge
def s_challenge():
    s = prs.slides.add_slide(L["Title only"])
    header(s, "The challenge", "Safety risk concentrates where road data is thinnest")

    pts = [
        "Disengagements and near misses cluster in uncommon mixes of weather, light, "
        "and road-user behaviour.",
        "A fleet can drive for years before it meets fog at dusk with a pedestrian "
        "stepping out from behind a vehicle.",
        "SOTIF (ISO 21448) expects evidence about the unsafe and unknown space, not "
        "only the miles already driven.",
    ]
    x, w = ML, 6.6
    y = TOP + 0.15
    for p in pts:
        oval(s, x + 0.02, y + 0.08, 0.13, 0.13, fill=BLUE)
        tf = tb(s, x + 0.34, y, w - 0.34, 1.1)
        para(tf, p, size=13.5, color=INK, line=1.2, after=0, first=True)
        y += 1.12

    tf = tb(s, x, y + 0.1, w, 0.6)
    para(tf, "The gap is not volume. It is the rare combinations that decide safety.",
         size=13, color=BLUE, bold=True, font=FONT_SB, line=1.15, after=0, first=True)

    gx, gy = 8.25, TOP + 0.5
    cols, rowsn = 6, 5
    cell, gap = 0.52, 0.12
    recorded = {(0, 0), (2, 1), (1, 3), (4, 0), (3, 4)}
    panel_w = cols * cell + (cols - 1) * gap
    tfh = tb(s, gx, gy - 0.42, panel_w, 0.3)
    para(tfh, "SCENARIO SPACE", size=10, color=SUB, bold=True, font=FONT_SB, spc=2, after=0, first=True)
    for r in range(rowsn):
        for c in range(cols):
            col = BLUE if (c, r) in recorded else GRID_OFF
            rect(s, gx + c * (cell + gap), gy + r * (cell + gap), cell, cell, fill=col)
    ly = gy + rowsn * (cell + gap) + 0.2
    rect(s, gx, ly + 0.02, 0.22, 0.22, fill=BLUE)
    para(tb(s, gx + 0.32, ly, 3.0, 0.3), "Recorded by the fleet", size=11, color=INK, after=0, first=True)
    rect(s, gx, ly + 0.36, 0.22, 0.22, fill=GRID_OFF)
    para(tb(s, gx + 0.32, ly + 0.34, 3.2, 0.3), "Still needs evidence", size=11, color=SUB, after=0, first=True)


# ================================================================== Slide 4: the approach
def s_approach():
    s = prs.slides.add_slide(L["Title only"])
    header(s, "The approach", "One pipeline generates scenarios and verdicts")

    steps = [
        ("1", "Seed", "Start from a small set of real driving clips."),
        ("2", "Generate", "Cosmos Transfer and Cosmos Predict create weather, lighting, "
                          "actor and behaviour variants."),
        ("3", "Validate", "The perception model runs on each clip; Cosmos Reason scores it "
                          "and explains every failure."),
        ("4", "Report", "Coverage view, failure gallery, and a one-page validation record."),
    ]
    n = len(steps)
    gap = 0.5
    cw = (CW - (n - 1) * gap) / n
    y, ch = TOP + 0.1, 2.55
    x = ML
    centers = []
    for num, name, body in steps:
        rect(s, x, y, cw, ch, fill=WHITE, line=HAIR, lw=1.0)
        rect(s, x, y, cw, 0.07, fill=BLUE)
        oval(s, x + 0.26, y + 0.32, 0.5, 0.5, fill=BLUE)
        tfn = tb(s, x + 0.26, y + 0.32, 0.5, 0.5, anchor=MSO_ANCHOR.MIDDLE)
        para(tfn, num, size=16, color=WHITE, bold=True, font=FONT_SB, align=PP_ALIGN.CENTER, after=0, first=True)
        tf = tb(s, x + 0.26, y + 1.0, cw - 0.52, ch - 1.15)
        para(tf, name, size=15, color=INK, bold=True, after=5, first=True)
        para(tf, body, size=11.5, color=SUB, line=1.16, after=0)
        centers.append((x + cw, y + ch / 2))
        x += cw + gap
    for cxp, cyp in centers[:-1]:
        chevron(s, cxp + gap / 2, cyp)

    ly = y + ch + 0.34
    rect(s, ML, ly, CW, 0.5, fill=BLUE_T)
    tf = tb(s, ML + 0.25, ly, CW - 0.5, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [("Closed loop.  ", True, BLUE, FONT_SB),
              ("New failures feed back into the scenario matrix, so coverage grows with each run.", False, INK, FONT)],
         size=12.5, after=0, first=True)

    ny = ly + 0.76
    tfn = tb(s, ML, ny, CW, 0.3)
    para(tfn, "BUILT ON THE NVIDIA PLATFORM", size=10, color=SUB, bold=True, font=FONT_SB, spc=2, after=0, first=True)
    comps = ["Cosmos Predict 2.5", "Cosmos Transfer 2.5", "Cosmos Reason",
             "NVIDIA AI Enterprise", "NGC containers", "H100 and A100 on AWS"]
    cx = ML
    cyp = ny + 0.36
    for comp in comps:
        chip_w = 0.22 + len(comp) * 0.082
        rrect(s, cx, cyp, chip_w, 0.34, fill=WHITE, line=HAIR, lw=1.0, radius=0.5)
        tf = tb(s, cx, cyp, chip_w, 0.34, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, comp, size=10.5, color=INK, align=PP_ALIGN.CENTER, after=0, first=True)
        cx += chip_w + 0.18


# ================================================================== Slide 5: evidence
def s_evidence():
    s = prs.slides.add_slide(L["Title only"])
    header(s, "Evidence from the proof of concept", "The first run already surfaced two unsafe responses")

    cards = [
        (YOLO_FOG, "Fog at dusk, occluded pedestrian", 3, R3,
         [("Detected late at frame 40, present from frame 29. The vehicle braked, but the "
           "margin was reduced.", INK, False),
          ("Cosmos Reason: late detection, action safe.", SUB, True)]),
        (YOLO_RAIN, "Rain at night, stalled vehicle", 5, R5,
         [("The stopped vehicle was not detected and the vehicle held its speed into "
           "braking traffic.", INK, False),
          ("Cosmos Reason: missed hazard, unsafe action.", SUB, True)]),
        (YOLO_SUN, "Low sun glare, cyclist cut-in", 4, R4,
         [("The cyclist was flagged late under glare. Braking began, but not early enough "
           "to clear the path.", INK, False),
          ("Cosmos Reason: late detection, unsafe action.", SUB, True)]),
    ]
    n = len(cards)
    gap = 0.4
    cw = (CW - (n - 1) * gap) / n
    y, ch = TOP + 0.05, 4.55
    x = ML
    for img, label, risk, rcol, body in cards:
        rect(s, x, y, cw, ch, fill=WHITE, line=HAIR, lw=1.0)
        pad = 0.18
        il, it, iw, ih = picture(s, img, x + pad, y + pad, cw - 2 * pad, 1.98)
        risk_pill(s, risk, rcol, il + iw, it)
        ty = y + pad + 1.98 + 0.16
        tfl = tb(s, x + pad, ty, cw - 2 * pad, 0.6)
        para(tfl, label, size=13.5, color=BLUE, bold=True, font=FONT_SB, line=1.05, after=0, first=True)
        hline(s, x + pad, x + cw - pad, ty + 0.62)
        tfb = tb(s, x + pad, ty + 0.74, cw - 2 * pad, ch - (ty - y) - 0.9)
        for i, (txt, col, ital) in enumerate(body):
            p = para(tfb, txt, size=11, color=col, line=1.18, after=6, first=(i == 0))
            if ital:
                for r in p.runs:
                    r.font.italic = True
        x += cw + gap

    tfn = tb(s, ML, y + ch + 0.14, CW, 0.4)
    para(tfn, "System under test: YOLOv8 perception. Judge: Cosmos Reason. Representative "
              "results from an initial run.", size=10, color=SUB, line=1.1, after=0, first=True)


# ================================================================== Slide 6: why MHP
def s_value():
    s = prs.slides.add_slide(L["Title only"])
    header(s, "Why MHP and what comes next", "MHP turns this into repeatable, audited validation")

    val = [
        ("Coverage by design", "We translate your hazard catalogue into a scenario matrix, so "
                               "testing targets the conditions that matter to your programme."),
        ("Verdicts you can defend", "Every clip carries a written rationale and a risk score, "
                                    "ready to drop into a SOTIF or ISO 26262 file."),
        ("Engineered for cost control", "Short GPU batches on AWS keep spend predictable while "
                                        "coverage grows with each run."),
    ]
    x, w = ML, 6.75
    y = TOP + 0.15
    for head_t, body in val:
        rect(s, x, y + 0.05, 0.045, 0.78, fill=BLUE)
        tf = tb(s, x + 0.24, y, w - 0.24, 1.1)
        para(tf, head_t, size=14.5, color=INK, bold=True, after=3, first=True)
        para(tf, body, size=12.5, color=SUB, line=1.18, after=0)
        y += 1.18

    px = 8.0
    pw = CR - px
    rect(s, px, TOP + 0.15, pw, 3.7, fill=BLUE_T)
    rect(s, px, TOP + 0.15, pw, 0.06, fill=BLUE)
    tf = tb(s, px + 0.3, TOP + 0.45, pw - 0.6, 0.4)
    para(tf, "NEXT STEPS", size=10.5, color=BLUE, bold=True, font=FONT_SB, spc=2, after=0, first=True)
    steps = [
        "Agree the priority hazards and the seed clips.",
        "Run a scoped pilot on your perception stack.",
        "Review the findings and a sample validation record together.",
    ]
    sy = TOP + 1.0
    for i, st in enumerate(steps, 1):
        oval(s, px + 0.3, sy, 0.34, 0.34, fill=BLUE)
        tfn = tb(s, px + 0.3, sy, 0.34, 0.34, anchor=MSO_ANCHOR.MIDDLE)
        para(tfn, str(i), size=12, color=WHITE, bold=True, font=FONT_SB, align=PP_ALIGN.CENTER, after=0, first=True)
        tfs = tb(s, px + 0.78, sy - 0.04, pw - 1.08, 0.85)
        para(tfs, st, size=12, color=INK, line=1.14, after=0, first=True)
        sy += 0.9

    by = 6.18
    rect(s, ML, by, CW, 0.66, fill=BLUE)
    tf = tb(s, ML + 0.3, by, CW - 0.6, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "From a handful of clips to a defensible safety story, without a new data campaign.",
         size=14, color=WHITE, bold=True, font=FONT_SB, after=0, first=True)


# ------------------------------------------------------------------ build
delete_all_slides()
s_cover()
s_exec()
s_challenge()
s_approach()
s_evidence()
s_value()
prs.save(str(OUT))
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
