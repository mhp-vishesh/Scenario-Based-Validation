"""Build the boss-facing OEM/client deck on the MHP corporate template.

Cover + 3 content slides (4 total), 16:9, MHP brand inherited from the template.

Output: Scenario-Based_Validation_OEM_Deck.pptx (workspace root)
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "MHP Company Profile EN.pptx"
DIAGRAM = ROOT / "assets" / "poc_plan_diagram.png"
SUT = ROOT / "POC" / "outputs" / "sut"
OUT = ROOT / "Scenario-Based_Validation_OEM_Deck.pptx"

YOLO_FOG = SUT / "20151221120048-D6-AGGRESSIVE-MOTORWAY_fog_dusk_pedestrian_occluded_emergence_000_yolo.jpg"
YOLO_RAIN = SUT / "20151221120048-D6-AGGRESSIVE-MOTORWAY_rain_night_stalled_vehicle_sudden_braking_001_yolo.jpg"
YOLO_SUN = SUT / "20151221120048-D6-AGGRESSIVE-MOTORWAY_clear_low_sun_glare_cyclist_jaywalk_cut_in_002_yolo.jpg"

# MHP palette (from the template theme "MHP colors")
MHP_BLUE = RGBColor(0x00, 0x00, 0x99)
LBLUE = RGBColor(0xCD, 0xD6, 0xFF)
GREEN = RGBColor(0x00, 0xCC, 0x67)
INK = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x57, 0x57, 0x57)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_FILL = RGBColor(0xF4, 0xF6, 0xFF)
# risk severity (safety convention)
R5 = RGBColor(0xC0, 0x39, 0x2B)
R4 = RGBColor(0xD3, 0x54, 0x00)
R3 = RGBColor(0xB5, 0x85, 0x2A)

FONT = "Segoe UI"


# ---------------------------------------------------------------- helpers
def layouts(prs):
    return {l.name: l for l in prs.slide_layouts}


def delete_all_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(idx)


def remove_ph(slide, idx):
    ph = get_ph(slide, idx)
    geom = (ph.left, ph.top, ph.width, ph.height)
    ph._element.getparent().remove(ph._element)
    return geom


def set_text(ph, text, size=None, bold=None, color=None, align=None):
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.name = FONT
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def add_text(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT,
         space_before=0, space_after=4, bullet=False, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bullet:
        _set_bullet(p)
    return p


def _set_bullet(p):
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial", "pitchFamily": "34", "charset": "0"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
    pPr.set("indent", str(-Inches(0.20)))
    pPr.set("marL", str(Inches(0.20)))
    pPr.append(buFont)
    pPr.append(buChar)


def contain_box(path, bl, bt, bw, bh):
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = bw / bh
    if ar > box_ar:
        w, h = bw, bw / ar
    else:
        h, w = bh, bh * ar
    return bl + (bw - w) / 2, bt + (bh - h) / 2, w, h


def add_picture_contain(slide, path, bl, bt, bw, bh):
    l, t, w, h = contain_box(str(path), bl, bt, bw, bh)
    pic = slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
    pic.line.color.rgb = LBLUE
    pic.line.width = Pt(0.75)
    return l, t, w, h


def risk_pill(slide, score, color, anchor_left_in, anchor_top_in):
    w, h = 1.02, 0.30
    left = anchor_left_in - w - 0.06
    top = anchor_top_in + 0.06
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"Risk {score}/5"
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE


# ---------------------------------------------------------------- build
prs = Presentation(str(TEMPLATE))
delete_all_slides(prs)
L = layouts(prs)
PAGE_W = Emu(prs.slide_width).inches

# --- Slide 1: cover -------------------------------------------------
s1 = prs.slides.add_slide(L["Title Slide"])
set_text(get_ph(s1, 0), "Scenario-Based Validation for ADAS", size=40, bold=True)
# the layout subtitle is a narrow lime "category" label; replace with a clean line
remove_ph(s1, 1)
sub = add_text(s1, 0.47, 4.05, 11.6, 1.3)
para(sub, "Closed-loop synthetic scenario generation and automated safety "
          "validation, built on NVIDIA Cosmos and AWS",
     18, color=WHITE, space_after=0, first=True)

# --- Slide 2: what it does + workflow + NVIDIA ----------------------
s2 = prs.slides.add_slide(L["Title only"])
t2 = get_ph(s2, 0)
t2.left, t2.top, t2.width, t2.height = Inches(0.45), Inches(0.45), Inches(12.43), Inches(0.85)
set_text(t2, "What the POC does and how it works", bold=True)

# left column: intro + NVIDIA platform
lc = add_text(s2, 0.45, 1.55, 4.95, 5.4)
para(lc, "From a small set of seed clips, the POC generates hundreds of edge-case "
         "variants and validates a perception model against every one. No extra "
         "fleet data collection.", 13.5, color=INK, space_after=12, first=True)
para(lc, "Built on the NVIDIA platform", 13.5, bold=True, color=MHP_BLUE, space_before=2, space_after=6)
para(lc, "Cosmos Predict 2.5 and Cosmos Transfer 2.5 generate the scenario variants.",
     12, color=INK, space_after=6, bullet=True)
para(lc, "Cosmos Reason reviews every clip and explains each failure.",
     12, color=INK, space_after=6, bullet=True)
para(lc, "NVIDIA AI Enterprise and NGC containers run on H100 and A100 GPUs (AWS).",
     12, color=INK, space_after=6, bullet=True)

# right column: workflow diagram
add_picture_contain(s2, DIAGRAM, 5.65, 1.5, 7.35, 5.2)

# --- Slide 3: generated clips + YOLO + verdict ----------------------
s3 = prs.slides.add_slide(L["3 images white"])
set_text(get_ph(s3, 18), "Representative POC results", size=11, color=GRAY)
set_text(get_ph(s3, 0), "Generated scenarios with automated safety verdicts", bold=True)

cols = [
    {
        "head": "Fog and dusk: pedestrian emerges from occlusion",
        "img": YOLO_FOG, "pic_idx": 21, "head_idx": 20, "body_idx": 1,
        "risk": 3, "risk_color": R3,
        "body": [
            ("YOLO detected the pedestrian at frame 40 (present from frame 29). "
             "The ego braked, but the safety margin was reduced.", INK, False),
            ("Cosmos Reason: late detection, action safe.", GRAY, True),
        ],
    },
    {
        "head": "Rain and night: stalled vehicle blocking the lane",
        "img": YOLO_RAIN, "pic_idx": 25, "head_idx": 22, "body_idx": 16,
        "risk": 5, "risk_color": R5,
        "body": [
            ("YOLO missed the stopped vehicle (no hazards flagged) and the ego "
             "kept its speed into braking traffic.", INK, False),
            ("Cosmos Reason: missed hazard, unsafe action.", GRAY, True),
        ],
    },
    {
        "head": "Low sun glare: cyclist cuts across the lane",
        "img": YOLO_SUN, "pic_idx": 26, "head_idx": 24, "body_idx": 17,
        "risk": 4, "risk_color": R4,
        "body": [
            ("YOLO flagged the cyclist late under glare. Braking started but not "
             "early enough to clear the path.", INK, False),
            ("Cosmos Reason: late detection, unsafe action.", GRAY, True),
        ],
    },
]

for c in cols:
    # headline placeholders carry a blue fill from the layout; use white text
    set_text(get_ph(s3, c["head_idx"]), c["head"], size=12.5, bold=True, color=WHITE)
    geom = remove_ph(s3, c["pic_idx"])
    bl, bt, bw, bh = (Emu(geom[0]).inches, Emu(geom[1]).inches,
                      Emu(geom[2]).inches, Emu(geom[3]).inches)
    il, it, iw, ih = add_picture_contain(s3, c["img"], bl, bt, bw, bh)
    risk_pill(s3, c["risk"], c["risk_color"], il + iw, it)
    body = get_ph(s3, c["body_idx"])
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, (txt, col, italic) in enumerate(c["body"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = txt
        r.font.name = FONT
        r.font.size = Pt(10.5)
        r.font.color.rgb = col
        r.font.italic = italic

# --- Slide 4: how MHP makes the difference --------------------------
s4 = prs.slides.add_slide(L["Title only"])
t4 = get_ph(s4, 0)
t4.left, t4.top, t4.width, t4.height = Inches(0.45), Inches(0.45), Inches(12.43), Inches(0.85)
set_text(t4, "How MHP makes the difference", bold=True)

cards = [
    ("01", "Coverage without new data collection",
     "We synthesize rare edge cases (weather, lighting, occlusions, actor "
     "behaviour) from a small seed set, so you test conditions the fleet has "
     "not recorded."),
    ("02", "Automated, explainable verdicts",
     "Cosmos Reason scores each clip and writes the reason for every failure, "
     "so findings are auditable and not just a pass or fail number."),
    ("03", "Standards-aligned evidence",
     "Results map to SOTIF (ISO 21448) and ISO 26262 needs, exported as a "
     "one-page validation report."),
]

card_w, gap, card_t, card_h = 3.94, 0.31, 1.65, 3.55
left = 0.45
for num, title, body in cards:
    rect = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(left), Inches(card_t), Inches(card_w), Inches(card_h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_FILL
    rect.line.color.rgb = LBLUE
    rect.line.width = Pt(1.0)
    rect.shadow.inherit = False
    rect.text_frame.word_wrap = True
    rect.text_frame.paragraphs[0].text = ""

    badge = s4.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(left + 0.28), Inches(card_t + 0.28), Inches(0.62), Inches(0.62))
    badge.fill.solid()
    badge.fill.fore_color.rgb = MHP_BLUE
    badge.line.fill.background()
    badge.shadow.inherit = False
    badge.text_frame.word_wrap = False
    bp = badge.text_frame.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = num
    br.font.name = FONT
    br.font.size = Pt(15)
    br.font.bold = True
    br.font.color.rgb = WHITE

    tf = add_text(s4, left + 0.30, card_t + 1.10, card_w - 0.60, card_h - 1.30)
    para(tf, title, 14, bold=True, color=MHP_BLUE, space_after=6, first=True)
    para(tf, body, 11.5, color=INK, space_after=0)
    left += card_w + gap

# takeaway bar
bar_t = card_t + card_h + 0.28
bar = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.45), Inches(bar_t), Inches(12.43), Inches(0.78))
bar.fill.solid()
bar.fill.fore_color.rgb = MHP_BLUE
bar.line.fill.background()
bar.shadow.inherit = False
bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
bp = bar.text_frame.paragraphs[0]
bp.alignment = PP_ALIGN.CENTER
br = bp.add_run()
br.text = ("We close the loop. Every new failure feeds back into the scenario "
           "matrix, so coverage grows with each run.")
br.font.name = FONT
br.font.size = Pt(13.5)
br.font.bold = True
br.font.color.rgb = WHITE

prs.save(str(OUT))
print("Saved:", OUT)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
